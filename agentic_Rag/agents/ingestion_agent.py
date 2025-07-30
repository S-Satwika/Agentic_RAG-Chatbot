import os
import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
import markdown
import textwrap

def chunk_text(text, max_tokens=300):
    """Simple chunker by paragraph or sentence."""
    return textwrap.wrap(text, max_tokens)

def read_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def read_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def read_txt_md(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def create_mcp_message(sender, receiver, payload, trace_id="trace-001", msg_type="PARSED_CONTENT"):
    return {
        "sender": sender,
        "receiver": receiver,
        "type": msg_type,
        "trace_id": trace_id,
        "payload": payload
    }

def process_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == ".pdf":
        raw_text = read_pdf(file_path)
    elif ext == ".docx":
        raw_text = read_docx(file_path)
    elif ext == ".pptx":
        raw_text = read_pptx(file_path)
    elif ext == ".csv":
        raw_text = read_csv(file_path)
    elif ext in [".txt", ".md"]:
        raw_text = read_txt_md(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    chunks = chunk_text(raw_text)
    return create_mcp_message(
        sender="IngestionAgent",
        receiver="RetrievalAgent",
        payload={
            "chunks": chunks,
            "doc_meta": {"filename": os.path.basename(file_path)}
        }
    )
